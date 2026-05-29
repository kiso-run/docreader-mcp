"""v0.2.0 — additional format support (PPTX, HTML, RTF, ODT, EPUB).

Markdown (.md) is already supported via the plain-text path in v0.1
since `.md` is in _TEXT_EXTENSIONS — covered here only as a regression
guard.

Each test builds its own minimal fixture in `tmp_path` to avoid
shipping binary fixtures for every format. The point is to verify the
reader handles the format without crash + extracts something
recognisable; deep parsing fidelity is the upstream library's job.
"""
from __future__ import annotations

from pathlib import Path

import pytest

from kiso_docreader_mcp.docreader_runner import (
    check_health,
    list_supported_formats,
    read_document,
)


# ── PPTX ────────────────────────────────────────────────────────


def _make_pptx(tmp_path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    blank_layout = prs.slide_layouts[5]  # title only
    slide1 = prs.slides.add_slide(blank_layout)
    slide1.shapes.title.text = "Quarterly review Q1 2026"
    slide2 = prs.slides.add_slide(blank_layout)
    slide2.shapes.title.text = "Revenue up 14% YoY"
    notes = slide2.notes_slide.notes_text_frame
    notes.text = "Speaker note: revenue driver was the new B2B tier"
    path = tmp_path / "sample.pptx"
    prs.save(str(path))
    return path


def test_pptx_extracts_slide_titles_and_speaker_notes(tmp_path: Path) -> None:
    path = _make_pptx(tmp_path)
    result = read_document(file_path=str(path))
    assert result["success"] is True
    assert result["format"] == "pptx"
    text = result["text"]
    assert "Quarterly review Q1 2026" in text
    assert "Revenue up 14% YoY" in text
    assert "revenue driver was the new B2B tier" in text


def test_pptx_emits_slide_headers(tmp_path: Path) -> None:
    path = _make_pptx(tmp_path)
    result = read_document(file_path=str(path))
    assert "# Slide 1" in result["text"]
    assert "# Slide 2" in result["text"]


# ── HTML ────────────────────────────────────────────────────────


def test_html_strips_script_style_and_keeps_body(tmp_path: Path) -> None:
    path = tmp_path / "page.html"
    path.write_text("""
        <html>
          <head>
            <script>alert('xss')</script>
            <style>body { color: red; }</style>
            <title>Page title</title>
          </head>
          <body>
            <nav><a href="/">Home</a></nav>
            <h1>About the company</h1>
            <p>We sell widgets that delight customers.</p>
            <footer>copyright 2026</footer>
          </body>
        </html>
    """)
    result = read_document(file_path=str(path))
    assert result["success"] is True
    assert result["format"] == "html"
    text = result["text"]
    assert "About the company" in text
    assert "widgets that delight" in text
    # The strip should drop script/style
    assert "alert(" not in text
    assert "color: red" not in text


# ── Markdown regression (already supported via text path) ───────


def test_markdown_still_read_via_text_path(tmp_path: Path) -> None:
    path = tmp_path / "notes.md"
    path.write_text("# Heading\n\nSome **bold** text and a [link](https://x).\n")
    result = read_document(file_path=str(path))
    assert result["success"] is True
    assert "Heading" in result["text"]
    assert "bold" in result["text"]


# ── RTF ─────────────────────────────────────────────────────────


def test_rtf_extracts_plain_text(tmp_path: Path) -> None:
    # Minimal valid RTF.
    rtf_source = r"""{\rtf1\ansi\deff0 {\fonttbl {\f0 Times New Roman;}}
\f0\fs24 This is a contract clause. The party agrees to pay {\b on time}.
}"""
    path = tmp_path / "doc.rtf"
    path.write_text(rtf_source)
    result = read_document(file_path=str(path))
    assert result["success"] is True
    assert result["format"] == "rtf"
    text = result["text"]
    assert "contract clause" in text
    assert "agrees to pay" in text


# ── ODT ─────────────────────────────────────────────────────────


def _make_odt(tmp_path: Path) -> Path:
    from odf.opendocument import OpenDocumentText
    from odf.text import P
    doc = OpenDocumentText()
    doc.text.addElement(P(text="OpenDocument sample paragraph 1"))
    doc.text.addElement(P(text="Second paragraph with policy info"))
    path = tmp_path / "sample.odt"
    doc.save(str(path))
    return path


def test_odt_extracts_paragraphs(tmp_path: Path) -> None:
    path = _make_odt(tmp_path)
    result = read_document(file_path=str(path))
    assert result["success"] is True
    assert result["format"] == "odt"
    text = result["text"]
    assert "OpenDocument sample paragraph 1" in text
    assert "policy info" in text


# ── EPUB ────────────────────────────────────────────────────────


def _make_epub(tmp_path: Path) -> Path:
    from ebooklib import epub
    book = epub.EpubBook()
    book.set_identifier("test-id-1")
    book.set_title("Test book")
    book.set_language("en")
    c1 = epub.EpubHtml(title="Chapter 1", file_name="chap_01.xhtml")
    c1.content = "<html><body><h1>Chapter 1</h1><p>Introduction text here</p></body></html>"
    c2 = epub.EpubHtml(title="Chapter 2", file_name="chap_02.xhtml")
    c2.content = "<html><body><h1>Chapter 2</h1><p>Second chapter content</p></body></html>"
    book.add_item(c1)
    book.add_item(c2)
    book.toc = (c1, c2)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", c1, c2]
    path = tmp_path / "book.epub"
    epub.write_epub(str(path), book)
    return path


def test_epub_extracts_chapter_text(tmp_path: Path) -> None:
    path = _make_epub(tmp_path)
    result = read_document(file_path=str(path))
    assert result["success"] is True
    assert result["format"] == "epub"
    text = result["text"]
    assert "Introduction text here" in text
    assert "Second chapter content" in text


# ── Catalog + health ───────────────────────────────────────────


def test_list_supported_formats_includes_new_set() -> None:
    out = list_supported_formats()
    structured = set(out["structured"])
    for ext in (".pdf", ".docx", ".xlsx", ".csv", ".pptx", ".html",
                ".rtf", ".odt", ".epub"):
        assert ext in structured, f"missing {ext}"


def test_check_health_validates_new_modules() -> None:
    result = check_health()
    # All new format libs must be importable for the runner image to
    # report healthy.
    assert result["healthy"] is True, result["issues"]
