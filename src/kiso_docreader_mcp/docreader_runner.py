"""Document text extraction core — PDF / DOCX / XLSX / CSV / plain text.

Each format has its own reader that extracts text and self-truncates at
semantic boundaries (page, row, paragraph) instead of at a mid-char chop.
Output budget is 50 000 chars; readers emit a `truncated` flag + a
continuation hint when exceeded.

All readers are pure-Python (pypdf, python-docx, openpyxl) — no external
binaries, no API keys.
"""
from __future__ import annotations

import csv
from pathlib import Path


_MAX_OUTPUT_CHARS = 50_000

_TEXT_EXTENSIONS = frozenset({
    ".txt", ".md", ".rst", ".log", ".json", ".yaml", ".yml",
    ".toml", ".ini", ".cfg", ".conf", ".sh", ".bash", ".py",
    ".js", ".ts", ".html", ".htm", ".xml", ".css", ".sql",
})

_STRUCTURED_EXTENSIONS = frozenset({
    ".pdf", ".docx", ".xlsx", ".csv",
    # v0.2.0
    ".pptx", ".html", ".htm", ".rtf", ".odt", ".epub",
})


def read_document(*, file_path: str, pages: str | None = None) -> dict:
    path = Path(file_path).expanduser()
    if not path.is_file():
        return _fail(f"file not found: {file_path}")

    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            return _read_pdf(path, pages)
        if ext == ".docx":
            return _read_docx(path)
        if ext == ".xlsx":
            return _read_xlsx(path)
        if ext == ".csv":
            return _read_csv(path)
        # v0.2.0 readers
        if ext == ".pptx":
            return _read_pptx(path)
        if ext in (".html", ".htm"):
            return _read_html(path)
        if ext == ".rtf":
            return _read_rtf(path)
        if ext == ".odt":
            return _read_odt(path)
        if ext == ".epub":
            return _read_epub(path)
        if ext in _TEXT_EXTENSIONS or _is_likely_text(path):
            return _read_text(path)
    except Exception as exc:
        return _fail(f"{ext or 'file'} extraction failed: {exc}")

    return _fail(
        f"unsupported format: {ext or '(no extension)'} — "
        "supported: pdf, docx, xlsx, csv, pptx, html, rtf, odt, epub, plain text"
    )


def document_info(*, file_path: str) -> dict:
    path = Path(file_path).expanduser()
    if not path.is_file():
        return {
            "success": False,
            "file_name": None,
            "size_bytes": None,
            "format": None,
            "pages": None,
            "sheets": None,
            "rows": None,
            "stderr": f"file not found: {file_path}",
        }

    ext = path.suffix.lower()
    info = {
        "success": True,
        "file_name": path.name,
        "size_bytes": path.stat().st_size,
        "format": ext.lstrip("."),
        "pages": None,
        "sheets": None,
        "rows": None,
        "stderr": "",
    }

    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            info["pages"] = len(PdfReader(str(path)).pages)
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(path), read_only=True)
            info["sheets"] = list(wb.sheetnames)
            wb.close()
        elif ext == ".csv":
            with open(path, newline="", encoding="utf-8", errors="replace") as f:
                info["rows"] = sum(1 for _ in csv.reader(f))
    except Exception as exc:
        info["success"] = False
        info["stderr"] = f"metadata extraction failed: {exc}"
    return info


def check_health() -> dict:
    issues: list[str] = []
    required_modules = (
        "pypdf", "docx", "openpyxl",
        # v0.2.0 additions
        "pptx", "bs4", "html2text", "striprtf", "odf", "ebooklib",
    )
    for mod in required_modules:
        try:
            __import__(mod)
        except ImportError:
            issues.append(f"Python module `{mod}` is not importable")
    return {"healthy": not issues, "issues": issues}


def list_supported_formats() -> dict:
    return {
        "structured": sorted(_STRUCTURED_EXTENSIONS),
        "text": sorted(_TEXT_EXTENSIONS),
        "notes": (
            "Unknown extensions are probed with a heuristic and read as "
            "plain text if the first 512 bytes look textual."
        ),
    }


def _read_pdf(path: Path, pages_spec: str | None) -> dict:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    total_pages = len(reader.pages)

    if pages_spec:
        indices = _parse_page_ranges(pages_spec, total_pages)
    else:
        indices = list(range(total_pages))

    parts: list[str] = []
    total_chars = 0
    truncated = False
    for i in indices:
        text = reader.pages[i].extract_text() or ""
        if not text.strip():
            continue
        page_text = f"--- Page {i + 1} ---\n{text.strip()}"
        if total_chars + len(page_text) > _MAX_OUTPUT_CHARS and parts:
            truncated = True
            break
        parts.append(page_text)
        total_chars += len(page_text)

    body = "\n\n".join(parts)
    return {
        "success": True,
        "text": body,
        "format": "pdf",
        "total_pages": total_pages,
        "pages_returned": len(parts),
        "truncated": truncated,
        "stderr": "" if parts else "no extractable text in the requested pages",
    }


def _read_docx(path: Path) -> dict:
    from docx import Document
    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    total_chars = sum(len(p) for p in paragraphs)

    kept: list[str] = []
    chars = 0
    for para in paragraphs:
        if chars + len(para) > _MAX_OUTPUT_CHARS and kept:
            break
        kept.append(para)
        chars += len(para)

    return {
        "success": True,
        "text": "\n\n".join(kept),
        "format": "docx",
        "total_chars": total_chars,
        "shown_chars": chars,
        "truncated": chars < total_chars,
        "stderr": "" if kept else "docx has no text content",
    }


def _read_xlsx(path: Path) -> dict:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    sheets = list(wb.sheetnames)

    parts: list[str] = []
    total_chars = 0
    truncated = False

    for sheet_name in sheets:
        if truncated:
            break
        ws = wb[sheet_name]
        rows: list[str] = []
        sheet_chars = 0
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if not any(cells):
                continue
            line = "\t".join(cells)
            if total_chars + sheet_chars + len(line) > _MAX_OUTPUT_CHARS and (parts or rows):
                truncated = True
                break
            rows.append(line)
            sheet_chars += len(line) + 1
        if rows:
            parts.append(
                f"--- Sheet: {sheet_name} ({len(rows)} rows) ---\n"
                + "\n".join(rows)
            )
            total_chars += sheet_chars

    wb.close()
    return {
        "success": True,
        "text": "\n\n".join(parts),
        "format": "xlsx",
        "sheets": sheets,
        "truncated": truncated,
        "stderr": "" if parts else "xlsx has no data",
    }


def _read_csv(path: Path) -> dict:
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        all_rows = list(reader)
    if not all_rows:
        return {
            "success": True,
            "text": "",
            "format": "csv",
            "total_rows": 0,
            "shown_rows": 0,
            "columns": [],
            "truncated": False,
            "stderr": "csv is empty",
        }

    columns = all_rows[0]
    lines: list[str] = []
    total_chars = 0
    for row in all_rows:
        line = "\t".join(row)
        if total_chars + len(line) > _MAX_OUTPUT_CHARS and lines:
            break
        lines.append(line)
        total_chars += len(line) + 1

    return {
        "success": True,
        "text": "\n".join(lines),
        "format": "csv",
        "total_rows": len(all_rows),
        "shown_rows": len(lines),
        "columns": columns,
        "truncated": len(lines) < len(all_rows),
        "stderr": "",
    }


def _read_text(path: Path) -> dict:
    text = path.read_text(encoding="utf-8", errors="replace")
    total_chars = len(text)
    if total_chars <= _MAX_OUTPUT_CHARS:
        return {
            "success": True,
            "text": text,
            "format": path.suffix.lower().lstrip(".") or "text",
            "total_chars": total_chars,
            "shown_chars": total_chars,
            "truncated": False,
            "stderr": "",
        }
    truncated = text[:_MAX_OUTPUT_CHARS]
    last_newline = truncated.rfind("\n")
    if last_newline > 0:
        truncated = truncated[:last_newline]
    return {
        "success": True,
        "text": truncated,
        "format": path.suffix.lower().lstrip(".") or "text",
        "total_chars": total_chars,
        "shown_chars": len(truncated),
        "truncated": True,
        "stderr": "",
    }


def _parse_page_ranges(spec: str, total: int) -> list[int]:
    indices: list[int] = []
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            start_s, end_s = part.split("-", 1)
            try:
                start_val = int(start_s)
                end_val = int(end_s)
            except ValueError as exc:
                raise ValueError(f"invalid page range: {part!r}") from exc
            if start_val > end_val:
                start_val, end_val = end_val, start_val
            start_i = max(0, start_val - 1)
            end_i = min(total, end_val)
            indices.extend(range(start_i, end_i))
        else:
            try:
                i = int(part) - 1
            except ValueError as exc:
                raise ValueError(f"invalid page number: {part!r}") from exc
            if 0 <= i < total:
                indices.append(i)
    return sorted(set(indices))


def _is_likely_text(path: Path) -> bool:
    try:
        sample = path.read_bytes()[:512]
    except OSError:
        return False
    if not sample:
        return False
    text_chars = sum(1 for b in sample if 32 <= b < 127 or b in (9, 10, 13))
    return text_chars / len(sample) > 0.85


def _fail(message: str) -> dict:
    return {
        "success": False,
        "text": "",
        "format": None,
        "truncated": False,
        "stderr": message,
    }


# ── v0.2.0 readers ──────────────────────────────────────────────


def _read_pptx(path: Path) -> dict:
    """Extract slide titles + body text + speaker notes from PPTX.

    Output: one section per slide with `# Slide N` header. Tables are
    serialised as tab-joined rows. Speaker notes get their own block.
    """
    from pptx import Presentation
    prs = Presentation(str(path))
    parts: list[str] = []
    total_chars = 0
    truncated = False
    total_slides = len(prs.slides)

    for idx, slide in enumerate(prs.slides, start=1):
        if truncated:
            break
        chunks: list[str] = [f"# Slide {idx}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        chunks.append(line)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        chunks.append("\t".join(cells))
        # Speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                chunks.append(f"[Speaker notes]\n{notes_text}")
        slide_text = "\n".join(chunks)
        if total_chars + len(slide_text) > _MAX_OUTPUT_CHARS and parts:
            truncated = True
            break
        parts.append(slide_text)
        total_chars += len(slide_text)

    body = "\n\n".join(parts)
    return {
        "success": True,
        "text": body,
        "format": "pptx",
        "total_slides": total_slides,
        "slides_returned": len(parts),
        "truncated": truncated,
        "stderr": "" if parts else "pptx has no text content",
    }


def _read_html(path: Path) -> dict:
    """Extract main content from HTML, stripping scripts/styles/nav/footer.

    Uses BeautifulSoup4 to drop noise, then html2text to render the
    remaining body as Markdown — preserves headings/lists/links/tables
    in an LLM-friendly form.
    """
    from bs4 import BeautifulSoup
    import html2text

    raw = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")
    for selector in ("script", "style", "nav", "footer", "noscript", "iframe"):
        for tag in soup.find_all(selector):
            tag.decompose()
    body = soup.body if soup.body is not None else soup

    converter = html2text.HTML2Text()
    converter.ignore_links = False
    converter.ignore_images = True
    converter.body_width = 0  # no hard wrap
    md_text = converter.handle(str(body)).strip()

    total_chars = len(md_text)
    if total_chars > _MAX_OUTPUT_CHARS:
        md_text = md_text[:_MAX_OUTPUT_CHARS]
        last_newline = md_text.rfind("\n")
        if last_newline > 0:
            md_text = md_text[:last_newline]
        truncated = True
    else:
        truncated = False

    return {
        "success": True,
        "text": md_text,
        "format": "html",
        "total_chars": total_chars,
        "shown_chars": len(md_text),
        "truncated": truncated,
        "stderr": "" if md_text else "html has no visible body content",
    }


def _read_rtf(path: Path) -> dict:
    """Extract plain text from RTF via the striprtf pure-Python parser."""
    from striprtf.striprtf import rtf_to_text

    raw = path.read_text(encoding="utf-8", errors="replace")
    text = rtf_to_text(raw).strip()
    total_chars = len(text)
    truncated = False
    if total_chars > _MAX_OUTPUT_CHARS:
        text = text[:_MAX_OUTPUT_CHARS]
        truncated = True
    return {
        "success": True,
        "text": text,
        "format": "rtf",
        "total_chars": total_chars,
        "shown_chars": len(text),
        "truncated": truncated,
        "stderr": "" if text else "rtf has no text content",
    }


def _read_odt(path: Path) -> dict:
    """Extract paragraph + heading text from ODT via odfpy."""
    from odf import teletype
    from odf.opendocument import load
    from odf.text import H, P

    doc = load(str(path))
    paragraphs: list[str] = []
    for elem in doc.getElementsByType(P) + doc.getElementsByType(H):
        line = teletype.extractText(elem).strip()
        if line:
            paragraphs.append(line)

    total_chars = sum(len(p) for p in paragraphs)
    kept: list[str] = []
    chars = 0
    for para in paragraphs:
        if chars + len(para) > _MAX_OUTPUT_CHARS and kept:
            break
        kept.append(para)
        chars += len(para)
    return {
        "success": True,
        "text": "\n\n".join(kept),
        "format": "odt",
        "total_chars": total_chars,
        "shown_chars": chars,
        "truncated": chars < total_chars,
        "stderr": "" if kept else "odt has no text content",
    }


def _read_epub(path: Path) -> dict:
    """Extract chapter text from EPUB.

    Walks ITEM_DOCUMENT entries (XHTML chapters) in spine order, runs
    each through BeautifulSoup for tag stripping, concats with chapter
    headers when titles are detectable.
    """
    from bs4 import BeautifulSoup
    import ebooklib
    from ebooklib import epub

    book = epub.read_epub(str(path))
    items = list(book.get_items_of_type(ebooklib.ITEM_DOCUMENT))
    parts: list[str] = []
    total_chars = 0
    truncated = False
    for item in items:
        if truncated:
            break
        soup = BeautifulSoup(item.get_body_content(), "html.parser")
        for selector in ("script", "style"):
            for tag in soup.find_all(selector):
                tag.decompose()
        chapter_text = soup.get_text("\n", strip=True)
        if not chapter_text:
            continue
        if total_chars + len(chapter_text) > _MAX_OUTPUT_CHARS and parts:
            truncated = True
            break
        parts.append(chapter_text)
        total_chars += len(chapter_text)

    return {
        "success": True,
        "text": "\n\n".join(parts),
        "format": "epub",
        "chapters_returned": len(parts),
        "truncated": truncated,
        "stderr": "" if parts else "epub has no readable chapters",
    }
